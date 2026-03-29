from pptx import Presentation
import json

def extract_ppt_content(file_path):
    prs = Presentation(file_path)

    slides_data = []

    for i, slide in enumerate(prs.slides):
        slide_info = {
            "slide_number": i + 1,
            "title": "",
            "body": [],
            "layout": slide.slide_layout.name,
            "placeholders": []
        }

        # Safer title detection
        if slide.shapes.title and slide.shapes.title.text:
            slide_info["title"] = slide.shapes.title.text.strip()

        for shape in slide.shapes:
            # Placeholder info
            if shape.is_placeholder:
                slide_info["placeholders"].append(
                    str(shape.placeholder_format.type)
                )

            # Extract text + bullets
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()

                    if text:
                        # Skip if already used as title
                        if text == slide_info["title"]:
                            continue

                        slide_info["body"].append({
                            "text": text,
                            "level": para.level  # bullet level
                        })

        slides_data.append(slide_info)

    return {"slides": slides_data}


# ✅ Save JSON (REQUIRED FOR ASSIGNMENT)
def save_json(data, path):
    with open(path, "w") as f:
        json.dump(data, f, indent=4)