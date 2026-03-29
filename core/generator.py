from pptx import Presentation
import os

def generate_presentation(data, template_path):
    prs = Presentation(template_path)

    for slide_data in data["slides"]:
        # Use title + content layout
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        # ✅ Title
        if slide.shapes.title and slide_data["title"]:
            slide.shapes.title.text = slide_data["title"]

        # ✅ Body
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()

        for item in slide_data["body"]:
            p = tf.add_paragraph()
            p.text = item["text"]      # ✅ FIXED
            p.level = item["level"]    # ✅ FIXED (bullet levels)

    # ✅ Ensure output folder exists
    os.makedirs("output", exist_ok=True)

    output_path = "output/final_presentation.pptx"
    prs.save(output_path)

    return output_path